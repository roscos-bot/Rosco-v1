"""Read Microsoft Office files — the Word/Excel/PowerPoint a business Drive is full
of, which Google's export can't touch (that only works on Google-native Docs). We
download the raw bytes and parse them locally: python-docx, openpyxl, python-pptx.

Each function returns '' if its library is missing or the file won't parse, so a
pull just degrades that file to 'unreadable' — it never crashes the pull. Legacy
binary .doc/.xls/.ppt are NOT handled here (they need a different toolchain); the
pull's type breakdown will flag them if any turn up.
"""
from __future__ import annotations

import io

# Office Open XML mime types (the modern .docx/.xlsx/.pptx).
DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def docx_text(data: bytes) -> str:
    try:
        import docx
        d = docx.Document(io.BytesIO(data))
        parts = [p.text for p in d.paragraphs if p.text and p.text.strip()]
        for tbl in d.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts).strip()
    except Exception:
        return ""


def xlsx_text(data: bytes, max_rows: int = 500) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        out = []
        for ws in wb.worksheets:
            out.append(f"# {ws.title}")
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows:
                    break
                cells = [str(c) for c in row if c not in (None, "")]
                if cells:
                    out.append(" | ".join(cells))
        wb.close()
        return "\n".join(out).strip()
    except Exception:
        return ""


def pptx_text(data: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        out = []
        for i, slide in enumerate(prs.slides, 1):
            out.append(f"# Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            out.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                        if cells:
                            out.append(" | ".join(cells))
        return "\n".join(out).strip()
    except Exception:
        return ""


def is_office(mime: str) -> bool:
    return (mime or "").lower() in (DOCX, XLSX, PPTX)


def office_text(data: bytes, mime: str) -> str:
    m = (mime or "").lower()
    if m == DOCX:
        return docx_text(data)
    if m == XLSX:
        return xlsx_text(data)
    if m == PPTX:
        return pptx_text(data)
    return ""
